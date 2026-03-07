"""Generate hackathon presentation slides."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TEAL = RGBColor(0x0D, 0x9B, 0x8C)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x6B, 0x70, 0x80)
LIGHT_BG = RGBColor(0xF0, 0xFD, 0xFA)
RED = RGBColor(0xDC, 0x26, 0x26)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_slide_content(slide, left, top, width, height, items, font_size=16, color=DARK):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
    return txBox


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ============================================================
# Slide 1: Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide, DARK)

add_text_box(slide, 1, 1.5, 11, 1.2, "Clinical Tutor", 54, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 2.8, 11, 0.8, "AI-Powered Clinical Teaching Assistant", 28, TEAL, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.0, 11, 0.6,
             "Training medical students on clinical reasoning with 2,000 real patient cases",
             18, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 1, 5.5, 11, 0.5,
             "AI / Data in Healthcare Hackathon  |  University of Toronto  |  March 2026",
             14, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 6.0, 11, 0.5, "Creative Track", 14, TEAL, True, PP_ALIGN.CENTER)

# ============================================================
# Slide 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 5, 0.7, "The Problem", 36, DARK, True)

add_bullet_slide_content(slide, 0.8, 1.5, 5.5, 4, [
    "Medical students need thousands of hours of",
    "case-based learning to develop clinical reasoning.",
    "",
    "But attending physicians have limited time",
    "for one-on-one bedside teaching.",
    "",
    "Existing alternatives:",
    "  \u2022  Static textbooks \u2014 no interactivity",
    "  \u2022  Generic AI chatbots \u2014 fabricate clinical data",
    "  \u2022  Simulation labs \u2014 expensive, limited cases",
], 18, DARK)

# Right side - key stat
txBox = add_text_box(slide, 7.5, 2.0, 5, 1.5, "2,000", 72, TEAL, True, PP_ALIGN.CENTER)
add_text_box(slide, 7.5, 3.5, 5, 0.6, "Real de-identified patient cases", 20, DARK, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.5, 4.1, 5, 0.6, "from MIMIC-III clinical database", 16, GRAY, False, PP_ALIGN.CENTER)

add_text_box(slide, 7.5, 5.2, 5, 0.5, "841K lab records  \u00b7  153K prescriptions", 16, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 7.5, 5.6, 5, 0.5, "23K diagnoses  \u00b7  6 structured tables", 16, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 3: Solution Overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Our Solution: Clinical Tutor", 36, DARK, True)

add_text_box(slide, 0.8, 1.4, 11, 0.6,
             "A 7-layer AI pipeline that teaches clinical reasoning through real patient cases.",
             20, GRAY)

# Three columns
cols = [
    ("Interactive Teaching", [
        "4-stage Socratic method",
        "Progressive case disclosure",
        "Adaptive to student answers",
        "Free Q&A after structured stages",
    ], 0.8),
    ("Grounded in Real Data", [
        "Knowledge Graph: 6,446 nodes",
        "Hybrid RAG retrieval",
        "LLM only synthesizes \u2014 never fabricates",
        "All data from MIMIC-III",
    ], 4.8),
    ("Data Quality as Teaching", [
        "Detects missing labs & contradictions",
        "Flags critical values automatically",
        "Turns data issues into learning moments",
        "Cross-validates summaries vs records",
    ], 8.8),
]

for title, items, left in cols:
    add_text_box(slide, left, 2.3, 3.8, 0.5, title, 20, TEAL, True)
    add_bullet_slide_content(slide, left, 3.0, 3.8, 3, [f"\u2022  {item}" for item in items], 15, DARK)

# ============================================================
# Slide 4: Architecture - 7 Layer Pipeline
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "7-Layer Pipeline Architecture", 36, DARK, True)

layers = [
    ("Layer 1", "Data Fusion", "data_loader.py", "6 CSV tables \u2192 2,000 Patient objects"),
    ("Layer 2", "Knowledge Graph", "knowledge_graph.py", "6,446 nodes, 220,520 edges (networkx)"),
    ("Layer 3", "Hybrid Retrieval", "rag_engine.py", "KG similarity (60%) + semantic embeddings (40%)"),
    ("Layer 4", "Data Quality", "data_quality.py", "Missing data, critical values, contradictions"),
    ("Layer 5", "LLM Reasoning", "reasoning.py", "Claude API, 4-stage Socratic teaching"),
    ("Layer 6", "API Server", "api.py", "FastAPI REST endpoints, session management"),
    ("Layer 7", "Frontend", "frontend/", "Next.js + TypeScript + Tailwind CSS"),
]

y_start = 1.5
for i, (label, name, file, desc) in enumerate(layers):
    y = y_start + i * 0.75
    # Layer number
    add_text_box(slide, 0.8, y, 1.2, 0.5, label, 14, WHITE, True, PP_ALIGN.CENTER)
    # Add bg rectangle for layer number
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(1.0), Inches(0.5)  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = label
    tf.paragraphs[0].font.size = Pt(12)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Layer name
    add_text_box(slide, 2.0, y, 2.5, 0.5, name, 16, DARK, True)
    # File
    add_text_box(slide, 4.5, y, 2.5, 0.5, file, 13, GRAY, False)
    # Description
    add_text_box(slide, 7.0, y, 5.5, 0.5, desc, 14, DARK)

# Bottom callout
add_text_box(slide, 0.8, 7.0, 11, 0.4,
             "\u26a0  The LLM is the last mile \u2014 it synthesizes real data from 4 upstream layers, not its own knowledge.",
             14, RED, True, PP_ALIGN.LEFT)

# ============================================================
# Slide 5: Knowledge Graph
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Knowledge Graph", 36, DARK, True)

add_text_box(slide, 0.8, 1.4, 5.5, 0.5, "4 Node Types, 4 Edge Types", 20, TEAL, True)

# Node types
add_bullet_slide_content(slide, 0.8, 2.2, 5, 2.5, [
    "\u2022  2,000 Patient nodes",
    "\u2022  2,461 Disease nodes (ICD-9 codes)",
    "\u2022  1,514 Drug nodes",
    "\u2022  471 Lab Test nodes",
    "",
    "Total: 6,446 nodes, 220,520 edges",
], 16, DARK)

# Edge types
add_text_box(slide, 7, 1.4, 5.5, 0.5, "Relationships", 20, TEAL, True)

add_bullet_slide_content(slide, 7, 2.2, 5.5, 2.5, [
    "Patient \u2500\u2500 HAS_DIAGNOSIS \u2500\u2500\u25b6 Disease",
    "Patient \u2500\u2500 TAKES_DRUG \u2500\u2500\u2500\u25b6 Drug",
    "Patient \u2500\u2500 HAS_LAB_TEST \u2500\u25b6 LabTest",
    "Disease \u2500\u2500 CO_OCCURS \u2500\u2500\u2500\u25b6 Disease",
    "",
    "Co-occurrence: pairs appearing in \u226510 patients",
], 16, DARK)

# Similar patient scoring
add_text_box(slide, 0.8, 5.0, 11, 0.5, "Similar Patient Scoring (Weighted Neighbor Overlap)", 20, TEAL, True)
add_bullet_slide_content(slide, 0.8, 5.7, 11, 1.5, [
    "Shared Diagnoses \u00d7 3.0    |    Shared Drugs \u00d7 1.5    |    Shared Lab Tests \u00d7 0.5",
    "No direct Patient\u2194Patient edges \u2014 similarity computed dynamically through shared neighbors",
], 16, DARK)

# ============================================================
# Slide 6: RAG Hybrid Retrieval
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Hybrid RAG Retrieval", 36, DARK, True)

add_text_box(slide, 0.8, 1.5, 5.5, 0.5, "Knowledge Graph Score (60%)", 22, TEAL, True)
add_bullet_slide_content(slide, 0.8, 2.2, 5.5, 2, [
    "\u2022  Structural similarity via shared neighbors",
    "\u2022  Weighted by clinical importance",
    "\u2022  Captures: same diseases, same drugs,",
    "    same lab panels",
], 16, DARK)

add_text_box(slide, 7, 1.5, 5.5, 0.5, "Semantic Score (40%)", 22, TEAL, True)
add_bullet_slide_content(slide, 7, 2.2, 5.5, 2, [
    "\u2022  sentence-transformers (all-MiniLM-L6-v2)",
    "\u2022  Embeddings stored in ChromaDB",
    "\u2022  Captures: narrative similarity in",
    "    discharge summaries",
], 16, DARK)

add_text_box(slide, 0.8, 4.5, 11, 0.6,
             "Final Score = KG Score \u00d7 0.6 + Semantic Score \u00d7 0.4",
             24, DARK, True, PP_ALIGN.CENTER)

add_text_box(slide, 0.8, 5.5, 11, 1, (
    "Why hybrid?  Two patients may have identical diagnoses and drugs (high KG score) "
    "but very different clinical narratives \u2014 or similar narratives but different structured data. "
    "Combining both signals gives clinically meaningful similarity."
), 16, GRAY, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 7: Data Quality as Teaching
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Data Quality \u2192 Teaching Moments", 36, DARK, True)

add_text_box(slide, 0.8, 1.4, 11, 0.5,
             "Real clinical data is messy. We turn that into an educational advantage.",
             18, GRAY)

checks = [
    ("Missing Lab Detection", "critical",
     "Discharge summary mentions troponin but no troponin in lab records \u2014 why might this happen?"),
    ("Critical Value Alerts", "critical",
     "Potassium 6.8 mEq/L (critical range: >6.0) \u2014 what's your immediate response?"),
    ("Drug-Diagnosis Consistency", "warning",
     "Patient on Warfarin but no documented coagulation disorder \u2014 what's the indication?"),
    ("Allergy Contradictions", "warning",
     "Penicillin allergy documented but Amoxicillin prescribed \u2014 is this safe?"),
]

y = 2.2
for title, severity, example in checks:
    sev_color = RED if severity == "critical" else RGBColor(0xD9, 0x77, 0x06)
    add_text_box(slide, 0.8, y, 3.5, 0.4, title, 16, sev_color, True)
    add_text_box(slide, 4.5, y, 8, 0.4, f'"{example}"', 14, GRAY)
    y += 0.7

add_text_box(slide, 0.8, 5.5, 11, 0.8,
             "Each finding becomes a Socratic question in the teaching session \u2014 "
             "the student must reason through the discrepancy before moving on.",
             16, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 8: 4-Stage Teaching Flow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "4-Stage Socratic Teaching Flow", 36, DARK, True)

stages = [
    ("Stage 1", "Case Presentation",
     "Demographics + chief complaint only.\nStudent builds initial differential diagnosis."),
    ("Stage 2", "Physical Examination",
     "Reveal vitals and exam findings.\nChallenge student: what changed in your differential?"),
    ("Stage 3", "Labs & Data Quality",
     "Present lab results + flag data quality issues.\nStudent proposes working diagnosis and treatment."),
    ("Stage 4", "Treatment & Comparison",
     "Reveal actual treatment. Compare with similar cases\nfrom Knowledge Graph. Discuss outcomes."),
]

y = 1.6
for num, title, desc in stages:
    # Stage number box
    shape = slide.shapes.add_shape(
        1, Inches(0.8), Inches(y), Inches(1.5), Inches(1.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = TEAL
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = title
    p2.font.size = Pt(11)
    p2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

    # Description
    add_text_box(slide, 2.6, y + 0.1, 5, 1, desc, 15, DARK)

    y += 1.3

# Free Q&A note
add_text_box(slide, 7.5, 2.0, 5, 1.5,
             "After Stage 4:", 18, TEAL, True)
add_text_box(slide, 7.5, 2.6, 5, 2, (
    "Free Q&A mode \u2014 student can ask\n"
    "anything about the case.\n\n"
    "LLM has full context from all\n"
    "4 previous stages + KG data +\n"
    "quality findings + similar cases."
), 15, DARK)

add_text_box(slide, 7.5, 5.0, 5, 0.8,
             "Cost per session: ~$0.01\nScalable to thousands of students",
             14, GRAY, False, PP_ALIGN.LEFT)

# ============================================================
# Slide 9: Reliability & Safety
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Reliability & Safety", 36, WHITE, True)

points = [
    ("Grounded Generation",
     "LLM only synthesizes data from upstream layers \u2014 cannot fabricate patient information"),
    ("Data Quality Awareness",
     "System actively detects and surfaces contradictions instead of hiding them"),
    ("Progressive Disclosure",
     "Information revealed step-by-step, preventing cognitive overload \u2014 mirrors real clinical workflow"),
    ("Honest Uncertainty",
     "When data is incomplete or ambiguous, the system says so explicitly"),
    ("Privacy by Design",
     "All data from MIMIC-III \u2014 de-identified with date-shifted timestamps, no real patient exposure"),
    ("Teaching, Not Deciding",
     "System teaches clinical reasoning \u2014 it never recommends treatment for actual patients"),
]

y = 1.6
for title, desc in points:
    add_text_box(slide, 0.8, y, 3.5, 0.4, "\u2713  " + title, 18, TEAL, True)
    add_text_box(slide, 4.5, y, 8, 0.4, desc, 15, GRAY)
    y += 0.8

# ============================================================
# Slide 10: Tech Stack
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, WHITE)

add_text_box(slide, 0.8, 0.5, 10, 0.7, "Tech Stack", 36, DARK, True)

stack = [
    ("Backend", "Python, FastAPI, networkx, ChromaDB, sentence-transformers, pandas"),
    ("Frontend", "Next.js 16, TypeScript, Tailwind CSS, React Markdown"),
    ("AI / LLM", "Claude API (Anthropic) \u2014 Haiku for dev, Sonnet for production"),
    ("Data", "MIMIC-III subset via HuggingFace (bavehackathon/2026-healthcare-ai)"),
    ("Embedding", "all-MiniLM-L6-v2 (384-dim sentence embeddings)"),
    ("Graph", "networkx \u2014 in-memory graph with weighted neighbor scoring"),
]

y = 1.5
for category, tech in stack:
    add_text_box(slide, 0.8, y, 2.5, 0.5, category, 18, TEAL, True)
    add_text_box(slide, 3.5, y, 9, 0.5, tech, 16, DARK)
    y += 0.7

# Stats
add_text_box(slide, 0.8, 5.8, 11, 0.5, "By the Numbers", 22, TEAL, True, PP_ALIGN.CENTER)
stats = "2,000 Cases  \u00b7  6,446 KG Nodes  \u00b7  220,520 KG Edges  \u00b7  841K Lab Records  \u00b7  153K Prescriptions"
add_text_box(slide, 0.8, 6.4, 11, 0.4, stats, 18, DARK, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 11: Live Demo
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_text_box(slide, 1, 2.5, 11, 1.2, "Live Demo", 60, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.0, 11, 0.6, "http://localhost:3000", 24, TEAL, False, PP_ALIGN.CENTER)

# ============================================================
# Slide 12: Thank You
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK)

add_text_box(slide, 1, 2.0, 11, 1, "Thank You", 54, WHITE, True, PP_ALIGN.CENTER)
add_text_box(slide, 1, 3.3, 11, 0.8, "Clinical Tutor", 28, TEAL, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 4.2, 11, 0.6,
             "AI-Powered Clinical Teaching Assistant", 20, GRAY, False, PP_ALIGN.CENTER)
add_text_box(slide, 1, 5.5, 11, 0.5, "Questions?", 24, WHITE, False, PP_ALIGN.CENTER)

# Save
output_path = r"e:\OpenGit\personal_project\healthcare_ai\Clinical_Tutor_Presentation.pptx"
prs.save(output_path)
print(f"Saved to: {output_path}")
